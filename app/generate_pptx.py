"""
PowerPoint Template Generator — RUL Copilot Presentation
=========================================================
Run: python app/generate_pptx.py

Generates a clean, professional PowerPoint template with placeholder slides
for the final presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
STEEL_BLUE = RGBColor(0x46, 0x82, 0xB4)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_RED = RGBColor(0xD3, 0x2F, 0x2F)
ACCENT_GREEN = RGBColor(0x38, 0x8E, 0x3C)
ACCENT_ORANGE = RGBColor(0xFF, 0xA0, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=DARK_GRAY):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_before = Pt(6)
        p.level = 0
    return txBox


def create_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # Accent bar
    add_shape(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), STEEL_BLUE)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Remaining Useful Life Prediction\nwith Uncertainty-Aware AI Copilot",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                 "NASA C-MAPSS Turbofan Engine Dataset",
                 font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "[Student Name]  •  [University]  •  [Date]",
                 font_size=16, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)


def create_agenda_slide(prs):
    """Slide 2: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                 "Agenda", font_size=28, bold=True, color=WHITE)

    items = [
        "1.  Problem Statement & Motivation",
        "2.  Dataset: NASA C-MAPSS FD001",
        "3.  Methodology & Model Architectures",
        "4.  Paper A Replication: CNN+LSTM + SG Smoothing",
        "5.  Paper B Replication: CNN-Transformer + MC Dropout",
        "6.  Novel Contribution: Uncertainty-Aware Copilot",
        "7.  Results & Ablation Study",
        "8.  Live Demo",
        "9.  Conclusion & Future Work",
    ]
    add_bullet_slide(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5), items, font_size=18)


def create_section_slide(prs, title, subtitle=""):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, Inches(4), Inches(3.4), Inches(5.3), Inches(0.04), STEEL_BLUE)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
                 title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                     subtitle, font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)


def create_content_slide(prs, title, bullets, note=""):
    """Standard content slide with header bar + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                 title, font_size=24, bold=True, color=WHITE)

    # Bullets
    add_bullet_slide(slide, Inches(1), Inches(1.6), Inches(11), Inches(4.5), bullets, font_size=17)

    # Footer note
    if note:
        add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                     note, font_size=12, color=RGBColor(0x88, 0x88, 0x88))

    return slide


def create_image_placeholder_slide(prs, title, caption="[Insert figure here]"):
    """Slide with placeholder for an image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                 title, font_size=24, bold=True, color=WHITE)

    # Placeholder box
    placeholder = add_shape(slide, Inches(2), Inches(1.8), Inches(9), Inches(4.5),
                            LIGHT_GRAY, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                 caption, font_size=16, color=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.CENTER)

    return slide


def create_table_slide(prs, title):
    """Slide with placeholder for a results table."""
    slide = create_content_slide(prs, title, [
        "Model                    MAE       RMSE      NASA Score",
        "─────────────────────────────────────────────────────",
        "Linear Regression         —          —           —",
        "Random Forest             —          —           —",
        "Gradient Boosting         —          —           —",
        "LSTM                      —          —           —",
        "CNN+LSTM (raw)            —          —           —",
        "CNN+LSTM (SG)             —          —           —",
        "CNN-Transformer           —          —           —",
        "CNN-Transformer + UQ      —          —           —",
    ], note="Fill in with actual experimental results.")
    return slide


def create_closing_slide(prs):
    """Thank you / Q&A slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
                 "Thank You", font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4), Inches(11.3), Inches(0.8),
                 "Questions?", font_size=24, color=LIGHT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.6),
                 "[student@email.com]  •  github.com/[username]",
                 font_size=14, color=RGBColor(0x99, 0xAA, 0xBB),
                 alignment=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Slide 1: Title ──
    create_title_slide(prs)

    # ── Slide 2: Agenda ──
    create_agenda_slide(prs)

    # ── Slide 3: Problem Statement ──
    create_section_slide(prs, "Problem Statement", "Why RUL Prediction Matters")

    create_content_slide(prs, "Motivation", [
        "Unplanned engine failures → safety risk + costly downtime",
        "Condition-based maintenance requires accurate RUL prediction",
        "Challenge: models must know when they DON'T know (uncertainty)",
        "Goal: build a system that predicts, explains, and knows its limits",
    ])

    # ── Slide 5: Dataset ──
    create_section_slide(prs, "Dataset", "NASA C-MAPSS FD001")

    create_content_slide(prs, "C-MAPSS FD001 Overview", [
        "Simulated turbofan engine run-to-failure data (NASA)",
        "100 training engines + 100 test engines",
        "1 operating condition, 1 fault mode (HPC degradation)",
        "21 sensor channels + 3 operational settings per cycle",
        "Target: predict Remaining Useful Life (cycles until failure)",
    ])

    create_image_placeholder_slide(prs, "Sensor Degradation Trends",
                                   "[Insert sensor trends plot from Week 1 EDA]")

    # ── Slide 8: Methodology ──
    create_section_slide(prs, "Methodology", "Preprocessing → Models → Uncertainty → Copilot")

    create_content_slide(prs, "Preprocessing Pipeline", [
        "Sensor selection: 14 informative sensors (7 near-constant dropped)",
        "Min-Max scaling per sensor across training set",
        "Optional: Savitzky-Golay smoothing (Paper A, window=11, poly=3)",
        "Sliding window: 30 cycles → predict RUL at window end",
        "RUL capping: min(raw_rul, 125) — piecewise linear",
        "Train/val split by engine unit (no data leakage)",
    ])

    create_content_slide(prs, "Model Architectures", [
        "Baseline: Linear Regression, Ridge, Random Forest, GBM",
        "LSTM: multi-layer, bidirectional option",
        "CNN+LSTM (Paper A): Conv1D → BatchNorm → ReLU → MaxPool → LSTM → FC",
        "CNN-Transformer (Paper B): Conv1D → Positional Encoding → Transformer Encoder → FC",
        "MC Dropout: N stochastic passes at inference → mean ± std",
    ])

    create_image_placeholder_slide(prs, "Architecture Diagram",
                                   "[Insert CNN-Transformer architecture diagram]")

    # ── Paper A ──
    create_section_slide(prs, "Paper A Replication",
                         "CNN+LSTM with Savitzky-Golay Smoothing")

    create_content_slide(prs, "Paper A: Key Findings", [
        "SG smoothing reduces noise while preserving degradation trends",
        "CNN extracts local patterns → LSTM captures temporal dependencies",
        "2×2 experiment: {LSTM, CNN+LSTM} × {raw, SG-smoothed}",
        "[Fill in: which combination performed best?]",
    ])

    create_image_placeholder_slide(prs, "SG Smoothing Effect",
                                   "[Insert raw vs SG-smoothed sensor comparison]")

    # ── Paper B ──
    create_section_slide(prs, "Paper B Replication",
                         "CNN-Transformer with MC Dropout UQ")

    create_content_slide(prs, "Paper B: Key Findings", [
        "Transformer self-attention captures long-range dependencies",
        "MC Dropout (N=100 passes) provides epistemic uncertainty estimates",
        "Higher uncertainty for engines closer to failure (as expected)",
        "[Fill in: calibration results, coverage at 95%]",
    ])

    create_image_placeholder_slide(prs, "Uncertainty Distribution",
                                   "[Insert MC Dropout distribution plot for sample engines]")

    # ── Novel Contribution ──
    create_section_slide(prs, "Novel Contribution",
                         "Uncertainty-Aware Ops Copilot with Abstention")

    create_content_slide(prs, "Copilot Pipeline", [
        "Step 1: Predict RUL + uncertainty (MC Dropout)",
        "Step 2: Explain prediction (top sensors via integrated gradients)",
        "Step 3: Retrieve knowledge base (BM25 over curated fault tree)",
        "Step 4: Generate recommendation with KB citation",
        "Safety: ABSTAIN when uncertainty > threshold → escalate to human",
    ])

    create_content_slide(prs, "Abstention Policy", [
        "Threshold sweep: plot MAE vs Coverage (% of accepted predictions)",
        "Key insight: refusing 10-15% of uncertain predictions can reduce MAE by X%",
        "OOD detection: Mahalanobis distance flags out-of-distribution inputs",
        "Decision matrix: {In-dist, OOD} × {Low unc, High unc}",
    ])

    create_image_placeholder_slide(prs, "Abstention Trade-off",
                                   "[Insert MAE vs Coverage plot]")

    create_image_placeholder_slide(prs, "Calibration Curve",
                                   "[Insert calibration curve: expected vs observed coverage]")

    # ── Results ──
    create_section_slide(prs, "Results", "Comparison & Ablations")

    create_table_slide(prs, "Model Comparison — FD001")

    create_content_slide(prs, "Ablation Study Highlights", [
        "Window size: 30 cycles optimal (vs 10, 20, 40, 50)",
        "SG smoothing: improves MAE by ~X% on average",
        "Attention heads: 4 heads = sweet spot (1 < 2 < 4 ≈ 8)",
        "MC samples: diminishing returns after ~50 passes",
    ])

    create_image_placeholder_slide(prs, "Model Comparison", "[Insert bar chart]")

    # ── Demo ──
    create_section_slide(prs, "Live Demo", "Streamlit Dashboard")

    create_content_slide(prs, "Demo Features", [
        "Single Engine Analysis: prediction + uncertainty + explanation",
        "Fleet Dashboard: all engines color-coded by risk",
        "Model Comparison: toggle between architectures",
        "Copilot Report: full predict → explain → retrieve → recommend",
    ])

    # ── Conclusion ──
    create_section_slide(prs, "Conclusion & Future Work")

    create_content_slide(prs, "Key Takeaways", [
        "Successfully replicated 2 recent papers (CNN+LSTM, CNN-Transformer)",
        "MC Dropout provides calibrated uncertainty estimates",
        "Abstention policy reduces error on accepted predictions",
        "Copilot agent grounds recommendations in KB + model explanation",
    ])

    create_content_slide(prs, "Future Work", [
        "Test on harder subsets: FD002 (6 operating conditions), FD004 (2 fault modes)",
        "Replace BM25 with dense retrieval (sentence-transformers)",
        "Add LLM-based copilot (GPT-4 / local model) for natural language",
        "Temporal attention visualization for deeper interpretability",
        "Deploy as REST API with FastAPI + Docker",
    ])

    # ── Closing ──
    create_closing_slide(prs)

    # ── Save ──
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, '..', 'reports', 'presentation_template.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {os.path.abspath(output_path)}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
