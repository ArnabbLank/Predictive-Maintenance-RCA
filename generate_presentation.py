"""
Generate Mid-Semester Presentation — RUL Copilot
Uses actual results from weeks 1–6 runs and generated plots.
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
STEEL_BLUE = RGBColor(0x46, 0x82, 0xB4)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_RED = RGBColor(0xD3, 0x2F, 0x2F)
ACCENT_GREEN = RGBColor(0x38, 0x8E, 0x3C)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ROOT = os.path.dirname(os.path.abspath(__file__))

def set_slide_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color: shape.line.color.rgb = line_color
    else: shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = 'Calibri'; p.space_before = Pt(6)
    return txBox

def header_bar(slide, title):
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                 title, font_size=24, bold=True, color=WHITE)

def section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, Inches(4), Inches(3.4), Inches(5.3), Inches(0.04), STEEL_BLUE)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
                 title, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                     subtitle, font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

def content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, title)
    add_bullets(slide, Inches(1), Inches(1.6), Inches(11), Inches(4.5), bullets, font_size=17)
    if note:
        add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                     note, font_size=12, color=RGBColor(0x88, 0x88, 0x88))
    return slide

def image_slide(prs, title, image_path, left=Inches(1.5), top=Inches(1.5),
                img_width=Inches(10), img_height=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, title)
    if os.path.exists(image_path):
        if img_height:
            slide.shapes.add_picture(image_path, left, top, img_width, img_height)
        else:
            slide.shapes.add_picture(image_path, left, top, width=img_width)
    else:
        add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                     f"[Image not found: {os.path.basename(image_path)}]",
                     font_size=16, color=RGBColor(0x99,0x99,0x99), alignment=PP_ALIGN.CENTER)
    return slide

def create_table(slide, left, top, width, height, headers, rows):
    """Add a table to a slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE
            p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = DARK_GRAY
                p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
            if r_idx % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape

def main():
    # Load results
    metrics_path = os.path.join(ROOT, 'reports', 'all_metrics.json')
    with open(metrics_path) as f:
        metrics = json.load(f)

    eda = metrics['eda']
    baselines = metrics['baselines']
    hi = metrics['health_index']
    lstm = metrics['lstm']
    paper_a = metrics['paper_a']
    ct = metrics['cnn_transformer']

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ═══ SLIDE 1: Title ═══
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.06), STEEL_BLUE)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Remaining Useful Life Prediction\nwith Uncertainty-Aware AI Copilot",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                 "Mid-Semester Progress Report — NASA C-MAPSS Turbofan Engine Dataset",
                 font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "Weeks 1–6 Results  •  EDA → Baselines → Deep Learning → Uncertainty",
                 font_size=16, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    # ═══ SLIDE 2: Agenda ═══
    content_slide(prs, "Agenda", [
        "1.  Problem Statement & Motivation",
        "2.  Dataset: NASA C-MAPSS FD001",
        "3.  Week 1: Exploratory Data Analysis",
        "4.  Week 2: ML Baselines (Linear, RF, GBM)",
        "5.  Week 3: Health Index Baseline",
        "6.  Week 4: LSTM Deep Learning Model",
        "7.  Week 5: Paper A — CNN+LSTM + SG Smoothing",
        "8.  Week 6: Paper B — CNN-Transformer + MC Dropout UQ",
        "9.  Results Comparison & Key Insights",
        "10. Agentic Copilot Demo & Next Steps",
    ])

    # ═══ SLIDE 3: Problem Statement ═══
    section_slide(prs, "Problem Statement", "Why RUL Prediction Matters")

    content_slide(prs, "Motivation", [
        "• Unplanned engine failures → safety risk + costly downtime",
        "• Condition-based maintenance requires accurate RUL prediction",
        "• Challenge: models must know WHEN they DON'T know (uncertainty)",
        "• Goal: build a system that predicts, explains, and knows its limits",
        "",
        "Our Approach:",
        "  Predict RUL → Quantify Uncertainty → Explain → Recommend → Abstain if unsure",
    ])

    # ═══ SLIDE 4-5: Dataset ═══
    section_slide(prs, "Dataset", "NASA C-MAPSS FD001")

    content_slide(prs, "C-MAPSS FD001 Overview", [
        f"• Simulated turbofan engine run-to-failure data (NASA)",
        f"• {eda['n_train_engines']} training engines + {eda['n_test_engines']} test engines",
        f"• 1 operating condition, 1 fault mode (HPC degradation)",
        f"• 21 sensor channels + 3 operational settings per cycle",
        f"• Engine lifetimes: {eda['min_lifetime']}–{eda['max_lifetime']} cycles (avg {eda['avg_lifetime']:.0f})",
        f"• Test RUL range: {eda['test_rul_min']}–{eda['test_rul_max']} (mean {eda['test_rul_mean']:.0f})",
        f"• RUL capping: min(raw_rul, 125) — piecewise linear target",
    ])

    # ═══ SLIDE 6: Week 1 EDA ═══
    section_slide(prs, "Week 1", "Exploratory Data Analysis")

    image_slide(prs, "Engine Lifetime Distribution",
                os.path.join(ROOT, 'week01-eda/temp/engine_lifetimes.png'),
                top=Inches(1.3), img_width=Inches(10.5))

    image_slide(prs, "RUL Label Distribution & Piecewise-Linear Target",
                os.path.join(ROOT, 'week01-eda/temp/rul_distributions.png'),
                top=Inches(1.3), img_width=Inches(11))

    image_slide(prs, "Sensor Variance Analysis — Identifying Informative Sensors",
                os.path.join(ROOT, 'week01-eda/temp/sensor_variance.png'),
                top=Inches(1.4), img_width=Inches(10))

    image_slide(prs, "Sensor Degradation Trends Over Engine Lifetime",
                os.path.join(ROOT, 'week01-eda/temp/sensor_trends.png'),
                top=Inches(1.3), img_width=Inches(10.5), img_height=Inches(5.5))

    image_slide(prs, "Sensor-RUL Correlation Heatmap",
                os.path.join(ROOT, 'week01-eda/temp/correlation_heatmap.png'),
                top=Inches(1.3), img_width=Inches(7), img_height=Inches(5.5))

    # Top correlated sensors
    top_sensors = list(eda['top_correlated_sensors'].keys())[:6]
    top_corrs = [f"{v:.3f}" for v in eda['top_correlated_sensors'].values()]
    content_slide(prs, "Week 1 Key Findings", [
        f"• 7 sensors have near-zero variance → dropped (sensor_1, 5, 6, 10, 16, 18, 19)",
        f"• 14 informative sensors retained for modeling",
        f"• Top correlated sensors with RUL: {', '.join(top_sensors[:3])}",
        f"  (|r| = {', '.join(top_corrs[:3])})",
        f"• Clear degradation trends visible in sensors over engine lifetime",
        f"• Decision: Window size=30, MinMax scaling, RUL cap=125",
    ])

    # ═══ SLIDE: Week 2 Baselines ═══
    section_slide(prs, "Week 2", "ML Baselines with Feature Engineering")

    # Results table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Baseline Model Results (FD001 Test)")
    bl_headers = ['Model', 'Test MAE', 'Test RMSE', 'NASA Score']
    bl_rows = [[b['Model'], f"{b['Test MAE']:.2f}", f"{b['Test RMSE']:.2f}", f"{b['Test NASA Score']:.0f}"]
               for b in baselines]
    create_table(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(3), bl_headers, bl_rows)
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.5),
                 "Features: rolling stats (mean, std, slope, min, max) + EWM + rate-of-change\n" +
                 "over last 30 cycles of each engine. Train/val split by engine unit.",
                 font_size=14, color=DARK_GRAY)

    image_slide(prs, "Baseline Predictions — Scatter Plots",
                os.path.join(ROOT, 'week02-baselines/temp/baseline_scatter.png'),
                top=Inches(1.3), img_width=Inches(10.5))

    image_slide(prs, "Random Forest Feature Importance",
                os.path.join(ROOT, 'week02-baselines/temp/feature_importance.png'),
                top=Inches(1.3), img_width=Inches(9), img_height=Inches(5.5))

    # ═══ SLIDE: Week 3 Health Index ═══
    section_slide(prs, "Week 3", "Health Index Baseline")

    image_slide(prs, "PCA Health Index — Per Engine & Relationship to RUL",
                os.path.join(ROOT, 'week03-health-index/temp/health_index.png'),
                top=Inches(1.3), img_width=Inches(11))

    content_slide(prs, "Health Index Results", [
        f"• PCA explained variance: 69.6%",
        f"• HI Baseline — MAE: {hi['Test MAE']:.2f}, RMSE: {hi['Test RMSE']:.2f}",
        f"• HI is a physics-inspired approach vs data-driven ML",
        f"• Takeaway: Single PCA component captures ~70% of variance",
        f"  but does not outperform ML baselines (single-point prediction from last cycle)",
        f"• ML baselines with window features provide much better RUL estimates",
    ])

    # ═══ SLIDE: Week 4 LSTM ═══
    section_slide(prs, "Week 4", "LSTM Deep Learning Model")

    image_slide(prs, "LSTM Training Curves",
                os.path.join(ROOT, 'week04-lstm-and-agent-skeleton/temp/lstm_training_curves.png'),
                top=Inches(1.3), img_width=Inches(10.5))

    image_slide(prs, "LSTM Predictions — Test Set",
                os.path.join(ROOT, 'week04-lstm-and-agent-skeleton/temp/lstm_scatter.png'),
                top=Inches(1.3), img_width=Inches(7))

    content_slide(prs, "LSTM Results & Agent Skeleton", [
        f"• LSTM (2 layers, hidden=64, dropout=0.3)",
        f"• Parameters: 57,985",
        f"• Test MAE: {lstm['MAE']:.2f}, RMSE: {lstm['RMSE']:.2f}, NASA: {lstm['NASA_Score']:.0f}",
        f"• Input: sliding windows of 30 cycles × 14 sensors",
        f"",
        f"• Also built Copilot agent skeleton:",
        f"  Tools: predict_rul, explain_prediction, retrieve_knowledge, recommend_action",
        f"  Orchestrator: rule-based pipeline (predict → explain → retrieve → recommend)",
    ])

    # ═══ SLIDE: Week 5 Paper A ═══
    section_slide(prs, "Week 5", "Paper A: CNN+LSTM + Savitzky-Golay Smoothing")

    image_slide(prs, "Savitzky-Golay Smoothing Effect on Sensor Signals",
                os.path.join(ROOT, 'week05-paper-a-cnn-lstm/temp/sg_smoothing_effect.png'),
                top=Inches(1.3), img_width=Inches(11))

    # Paper A results table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Paper A — 2×2 Experiment Results")
    pa_headers = ['Configuration', 'MAE', 'RMSE', 'NASA Score']
    pa_rows = [[r['name'], f"{r['MAE']:.2f}", f"{r['RMSE']:.2f}", f"{r['NASA_Score']:.0f}"]
               for r in paper_a]
    create_table(slide, Inches(2), Inches(1.8), Inches(9), Inches(3), pa_headers, pa_rows)
    add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(1.5),
                 "2×2 experiment: {LSTM, CNN+LSTM} × {raw, SG-smoothed}\n" +
                 "All models trained for 50 epochs with early stopping (patience=10)",
                 font_size=14, color=DARK_GRAY)

    image_slide(prs, "Paper A — MAE Comparison",
                os.path.join(ROOT, 'week05-paper-a-cnn-lstm/temp/papera_comparison.png'),
                top=Inches(1.3), img_width=Inches(9))

    # ═══ SLIDE: Week 6 Paper B ═══
    section_slide(prs, "Week 6", "Paper B: CNN-Transformer + MC Dropout UQ")

    content_slide(prs, "CNN-Transformer Architecture", [
        "• 1D-CNN feature extractor (2 conv layers → BatchNorm → ReLU)",
        "• Positional Encoding (sinusoidal)",
        "• Transformer Encoder (2 layers, 4 heads, d_model=64)",
        "• Global Average Pooling → FC regression head",
        "• MC Dropout (rate=0.1) kept active at inference",
        "• 84,417 trainable parameters",
        "",
        f"• Test MAE: {ct['ct_mae']:.2f}, RMSE: {ct['ct_rmse']:.2f}, NASA: {ct['ct_nasa']:.0f}",
    ])

    image_slide(prs, "CNN-Transformer Training Curves",
                os.path.join(ROOT, 'week06-paper-b-transformer-uq/temp/ct_training_curves.png'),
                top=Inches(1.3), img_width=Inches(10.5))

    image_slide(prs, "Predictions with 95% Confidence Intervals",
                os.path.join(ROOT, 'week06-paper-b-transformer-uq/temp/ct_predictions_with_ci.png'),
                top=Inches(1.3), img_width=Inches(11))

    image_slide(prs, "Uncertainty Analysis — Sanity Checks",
                os.path.join(ROOT, 'week06-paper-b-transformer-uq/temp/uncertainty_analysis.png'),
                top=Inches(1.3), img_width=Inches(11))

    image_slide(prs, "MC Dropout Prediction Distributions (Sample Engines)",
                os.path.join(ROOT, 'week06-paper-b-transformer-uq/temp/mc_distributions.png'),
                top=Inches(1.3), img_width=Inches(11))

    content_slide(prs, "Uncertainty Quantification — Key Findings", [
        f"• MC Dropout with N={100} stochastic forward passes",
        f"• Uncertainty-Error Correlation: {ct['uncertainty_error_corr']:.3f}",
        f"  → Positive correlation = model is more uncertain when it makes bigger errors",
        f"• Mean uncertainty (std): {ct['mean_uncertainty']:.2f} cycles",
        f"• Uncertainty enables:",
        f"  - Abstention: refuse prediction when uncertainty too high",
        f"  - Confidence intervals: communicate reliability to operators",
        f"  - Risk stratification: prioritize engines by uncertainty level",
    ])

    # ═══ SLIDE: Combined Results ═══
    section_slide(prs, "All Results", "Model Comparison — FD001 Test Set")

    # Master results table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Model Comparison — FD001 Test Set")
    all_headers = ['Model', 'MAE', 'RMSE', 'NASA Score']
    all_rows = [
        ['Gradient Boosting', f"{baselines[3]['Test MAE']:.2f}", f"{baselines[3]['Test RMSE']:.2f}", f"{baselines[3]['Test NASA Score']:.0f}"],
        ['Random Forest', f"{baselines[2]['Test MAE']:.2f}", f"{baselines[2]['Test RMSE']:.2f}", f"{baselines[2]['Test NASA Score']:.0f}"],
        ['Linear Regression', f"{baselines[0]['Test MAE']:.2f}", f"{baselines[0]['Test RMSE']:.2f}", f"{baselines[0]['Test NASA Score']:.0f}"],
        ['LSTM', f"{lstm['MAE']:.2f}", f"{lstm['RMSE']:.2f}", f"{lstm['NASA_Score']:.0f}"],
        ['LSTM (raw)', f"{paper_a[0]['MAE']:.2f}", f"{paper_a[0]['RMSE']:.2f}", f"{paper_a[0]['NASA_Score']:.0f}"],
        ['CNN-Transformer (MC)', f"{ct['ct_mae']:.2f}", f"{ct['ct_rmse']:.2f}", f"{ct['ct_nasa']:.0f}"],
    ]
    create_table(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(4.5), all_headers, all_rows)
    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                 "Lower MAE/RMSE/NASA Score = Better. NASA Score penalizes late predictions more heavily.",
                 font_size=12, color=RGBColor(0x88, 0x88, 0x88))

    image_slide(prs, "Model Comparison — Bar Charts",
                os.path.join(ROOT, 'reports/model_comparison.png'),
                top=Inches(1.3), img_width=Inches(11))

    # ═══ NEW SLIDES: Why CNN-Transformer over Gradient Boosting? ═══
    section_slide(prs, "Why CNN-Transformer?",
                  "Gradient Boosting Has Better Numbers — So Why Not Use It?")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Honest Comparison: Gradient Boosting vs CNN-Transformer")
    gb_mae = baselines[3]['Test MAE']
    gb_rmse = baselines[3]['Test RMSE']
    gb_nasa = baselines[3]['Test NASA Score']
    ct_mae_val = ct['ct_mae']
    ct_rmse_val = ct['ct_rmse']
    ct_nasa_val = ct['ct_nasa']
    mae_winner = 'GB' if gb_mae < ct_mae_val else 'CT ✓'
    rmse_winner = 'GB' if gb_rmse < ct_rmse_val else 'CT ✓'
    nasa_winner = 'GB' if gb_nasa < ct_nasa_val else 'CT ✓'
    mae_diff = abs(gb_mae - ct_mae_val)
    comp_headers = ['Metric', 'Gradient Boosting', 'CNN-Transformer (MC)', 'Winner']
    comp_rows = [
        ['MAE', f"{gb_mae:.2f}", f"{ct_mae_val:.2f}", f'{mae_winner} (Δ={mae_diff:.2f})'],
        ['RMSE', f"{gb_rmse:.2f}", f"{ct_rmse_val:.2f}", rmse_winner],
        ['NASA Score', f"{gb_nasa:.0f}", f"{ct_nasa_val:.0f}", nasa_winner],
        ['Uncertainty (std)', '✗ Not available', f"{ct['mean_uncertainty']:.2f} cycles", 'CT ✓'],
        ['Confidence Intervals', '✗ Not available', '95% CI per prediction', 'CT ✓'],
        ['Abstention / Escalation', '✗ Not possible', '✓ Knows when unsure', 'CT ✓'],
        ['Temporal Patterns', 'Summarized (flat)', 'Reads raw sequences', 'CT ✓'],
        ['Explainability (Grad)', '✗ Feature importance only', '✓ Per-timestep saliency', 'CT ✓'],
    ]
    create_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.0), comp_headers, comp_rows)
    accuracy_note = (f"GB MAE={gb_mae:.2f} vs CT MAE={ct_mae_val:.2f} (Δ={mae_diff:.2f}). "
                     f"GB NASA={gb_nasa:.0f} vs CT NASA={ct_nasa_val:.0f}. "
                     "CT provides uncertainty, explainability, and abstention that GB cannot.")
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
                 accuracy_note, font_size=12, color=RGBColor(0x88, 0x88, 0x88))

    content_slide(prs, "Why Not Gradient Boosting? — The Safety Argument", [
        "Gradient Boosting gives ONE number: 'RUL = 50 cycles.' That's it.",
        "A maintenance engineer must ask: How confident is the model? → No answer.",
        "",
        "CNN-Transformer (MC Dropout) gives:",
        "  • RUL = 50 cycles ± 5 cycles (95% CI: [40, 60]) — I'm confident",
        "  • RUL = 50 cycles ± 30 cycles (95% CI: [0, 110]) — DON'T TRUST ME",
        "",
        "In aviation, a wrong confident prediction is catastrophic.",
        "It's BETTER to say 'I don't know' than to say '50 cycles' and be wrong.",
        "",
        "The CNN-Transformer is the only model that can do this.",
        f"Uncertainty-error correlation: {ct['uncertainty_error_corr']:.3f} — the model IS more uncertain when wrong.",
    ])

    content_slide(prs, "How We Will Improve CNN-Transformer Results", [
        f"The CNN-Transformer's MAE ({ct['ct_mae']:.2f}) is close to GB ({baselines[3]['Test MAE']:.2f}).",
        "Planned improvements to beat GB on raw accuracy too:",
        "",
        "Week 10 — Abstention: Skip high-uncertainty predictions →",
        "  remaining predictions become MORE accurate (lower MAE on confident subset)",
        "",
        "Week 11 — Ablation & Tuning:",
        "  • Window size optimization (10/20/30/40/50)",
        "  • Attention heads (2/4/8) and model depth tuning",
        "  • Learning rate scheduling refinement",
        "  • Ensemble: average predictions from multiple MC Dropout runs",
        "",
        "Additional ideas: Feature scaling normalization, longer training,",
        "  curriculum learning (easy → hard samples), loss function tuning (Huber loss)",
    ])

    # ═══ NEW SLIDES: Copilot Architecture ═══
    section_slide(prs, "Copilot Architecture", "How the Ops Copilot Works — Current & Future")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Current Copilot: Rule-Based Pipeline (No LLM)")
    # Draw the pipeline diagram using text boxes and shapes
    steps = [
        ("1. PREDICT", "predict_rul()", "CNN-Transformer\n+ MC Dropout\n→ RUL mean, std,\n   95% CI", Inches(0.3)),
        ("2. EXPLAIN", "explain_prediction()", "Integrated Gradients\n→ Top sensors\n   ranked by\n   attribution", Inches(3.35)),
        ("3. RETRIEVE", "retrieve_knowledge()", "BM25 search over\nkb/fault_tree.md\n& glossary.md\n→ KB snippets", Inches(6.4)),
        ("4. RECOMMEND", "recommend_action()", "Rule-based logic:\nIF uncertainty > 20\n  → ESCALATE\nELSE → recommend", Inches(9.45)),
    ]
    for title, func, desc, left in steps:
        add_shape(slide, left, Inches(1.6), Inches(3.0), Inches(1.0), NAVY)
        add_text_box(slide, left + Inches(0.1), Inches(1.7), Inches(2.8), Inches(0.8),
                     title, font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_shape(slide, left, Inches(2.7), Inches(3.0), Inches(0.5), STEEL_BLUE)
        add_text_box(slide, left + Inches(0.1), Inches(2.75), Inches(2.8), Inches(0.4),
                     func, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_shape(slide, left, Inches(3.3), Inches(3.0), Inches(2.2), LIGHT_BLUE)
        add_text_box(slide, left + Inches(0.1), Inches(3.4), Inches(2.8), Inches(2.0),
                     desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow indicators between boxes
    for x_pos in [Inches(3.3), Inches(6.35), Inches(9.4)]:
        add_text_box(slide, x_pos, Inches(2.3), Inches(0.3), Inches(0.4),
                     "→", font_size=28, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1.0),
                 "Current status: 100% rule-based Python. No LLM, no API calls.\n"
                 "The orchestrator is a Python class (OpsCopilot) that calls 4 functions in fixed order with if/else logic.",
                 font_size=14, color=DARK_GRAY)

    content_slide(prs, "Safety Guardrails — Built Into Every Analysis", [
        "Before any recommendation is shown, the Copilot verifies 3 safety checks:",
        "",
        "✓ Uncertainty Reported — Never give advice without showing confidence level",
        "   If uncertainty is not available → flag as incomplete",
        "",
        "✓ Knowledge Base Cited — Recommendations must be grounded in the KB",
        "   If no relevant KB entry found → note 'No KB reference available'",
        "",
        "✓ No False Reassurance — If uncertainty is high, never say 'all clear'",
        "   High uncertainty + 'recommend' + 'high confidence' → BLOCKED",
        "",
        "Abstention Decision:",
        "  Uncertainty (std) > 20 cycles → ESCALATE to human expert",
        "  Uncertainty level = 'high' or 'very_high' → ESCALATE to human expert",
    ])

    # ═══ NEW SLIDES: Agentic Roadmap ═══
    section_slide(prs, "Agentic Roadmap", "From Rule-Based → LLM-Powered Copilot")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE); header_bar(slide, "Evolution of the Copilot — Week by Week")
    evo_headers = ['Week', 'What Changes', 'Agent Type']
    evo_rows = [
        ['Week 4', 'Agent skeleton: 4 tools defined, fixed pipeline', 'Rule-based'],
        ['Week 7', 'Explainability tools added (saliency, IG, SHAP)', 'Rule-based'],
        ['Week 8', 'Knowledge base + BM25 retrieval integrated', 'Rule-based'],
        ['Week 9', 'Full pipeline wired: single-engine + fleet analysis', 'Rule-based'],
        ['Week 10', 'Novelty detection + abstention logic formalized', 'Rule-based + OOD'],
        ['Week 11', 'Streamlit demo app with interactive controls', 'Rule-based (UI)'],
        ['Stretch', 'LLM orchestrator: OpenAI function calling via LangChain', 'LLM-backed'],
    ]
    create_table(slide, Inches(1), Inches(1.5), Inches(11), Inches(5.0), evo_headers, evo_rows)
    add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.5),
                 "Core tools stay the same. Only the orchestrator changes (if/else → LLM tool calling).",
                 font_size=13, color=RGBColor(0x88, 0x88, 0x88))

    content_slide(prs, "Final Product Vision — The Streamlit Dashboard", [
        "Interactive Streamlit app with 3 pages:",
        "",
        "1. Single Engine Analysis:",
        "   Select engine → Run 4-step Copilot → See prediction, uncertainty,",
        "   sensor attribution chart, KB reference, and recommendation / escalation",
        "",
        "2. Fleet Dashboard:",
        "   Analyze all 100 engines → Heatmap of risk levels → Sorted by urgency",
        "   Summary: X critical, Y warning, Z normal, W escalated",
        "",
        "3. Ablation & Comparison:",
        "   Compare all models side-by-side, toggle UQ on/off, vary MC samples",
        "",
        "The tools (predict, explain, retrieve, recommend) will stay rule-based.",
        "LLM integration is a stretch goal — the value is in uncertainty + explanations.",
    ])

    # ═══ SLIDE: Copilot & Next Steps ═══
    section_slide(prs, "Agentic Copilot", "Predict → Explain → Retrieve → Recommend")

    content_slide(prs, "Ops Copilot — Pipeline", [
        "Step 1: Predict RUL + uncertainty (MC Dropout, CNN-Transformer)",
        "Step 2: Explain prediction (gradient saliency / integrated gradients)",
        "Step 3: Retrieve knowledge base (BM25 over curated fault tree)",
        "Step 4: Generate recommendation with KB citation",
        "",
        "Safety Guardrails:",
        "  • ABSTAIN when uncertainty > threshold → escalate to human",
        "  • Verify uncertainty is reported, KB is cited",
        "  • No false reassurance when uncertainty is high",
    ])

    content_slide(prs, "Mid-Semester Summary & Key Insights", [
        f"✓ Comprehensive EDA on C-MAPSS FD001 — 14 informative sensors identified",
        f"✓ ML baselines: Gradient Boosting best at MAE={baselines[3]['Test MAE']:.2f}",
        f"✓ LSTM sequence model: MAE={lstm['MAE']:.2f}",
        f"✓ CNN-Transformer with MC Dropout: MAE={ct['ct_mae']:.2f} with uncertainty",
        f"✓ Uncertainty-error correlation: {ct['uncertainty_error_corr']:.3f} (model knows when unsure)",
        f"✓ Copilot agent skeleton built: predict → explain → retrieve → recommend",
    ])

    content_slide(prs, "Remaining Work (Weeks 7–12)", [
        "Week 7:  Explainability — saliency maps, integrated gradients, SHAP",
        "Week 8:  Knowledge Base RAG — BM25 + sentence-transformer retrieval",
        "Week 9:  Agentic Copilot — LLM integration with tool calling",
        "Week 10: Novelty detection & abstention policy tuning",
        "Week 11: Ablation studies & Streamlit dashboard",
        "Week 12: Final report & presentation",
    ])

    # ═══ SLIDE: Closing ═══
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
                 "Thank You", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4), Inches(11.3), Inches(0.8),
                 "Questions & Demo", font_size=24, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # ═══ Save ═══
    output_path = os.path.join(ROOT, 'reports', 'mid_semester_presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {os.path.abspath(output_path)}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
